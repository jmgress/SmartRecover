#!/usr/bin/env python3
"""
Generate a PowerPoint architecture diagram for SmartRecover.

This script produces docs/SmartRecover_Architecture.pptx — an editable one-slide
PowerPoint that shows the 3-panel UI layout and the LangGraph agent orchestration
flow so stakeholders can customise individual components.

Usage:
    pip install python-pptx
    python docs/generate_architecture_diagram.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ---------------------------------------------------------------------------
# Colours (matching the SmartRecover UI palette)
# ---------------------------------------------------------------------------
C_BG          = RGBColor(0xF8, 0xF9, 0xFC)   # slide background
C_PANEL_BG    = RGBColor(0xFF, 0xFF, 0xFF)   # panel fill
C_HEADER      = RGBColor(0x4C, 0x1D, 0x95)   # purple header
C_HEADER_TEXT = RGBColor(0xFF, 0xFF, 0xFF)   # white header text
C_BORDER      = RGBColor(0xE5, 0xE7, 0xEB)   # light grey border
C_TITLE_TEXT  = RGBColor(0x1F, 0x29, 0x37)   # dark text
C_SUBTEXT     = RGBColor(0x6B, 0x72, 0x80)   # muted text
C_AGENT       = RGBColor(0x7C, 0x3A, 0xED)   # agent purple
C_AGENT_LIGHT = RGBColor(0xED, 0xE9, 0xFE)   # agent card bg
C_ARROW       = RGBColor(0x9C, 0xA3, 0xAF)   # connector lines
C_CRIT        = RGBColor(0xDC, 0x26, 0x26)   # critical red
C_HIGH        = RGBColor(0xEA, 0x58, 0x0C)   # high orange
C_MED         = RGBColor(0xD9, 0x77, 0x06)   # medium yellow
C_LOW         = RGBColor(0x25, 0x63, 0xEB)   # low blue
C_GREEN       = RGBColor(0x16, 0x65, 0x34)   # resolved green
C_GREEN_BG    = RGBColor(0xDC, 0xFC, 0xE7)
C_YELLOW_BG   = RGBColor(0xFE, 0xF3, 0xC7)
C_RED_BG      = RGBColor(0xFE, 0xE2, 0xE2)
C_BLUE_BG     = RGBColor(0xDB, 0xEA, 0xFE)
C_ORANGE_BG   = RGBColor(0xFF, 0xED, 0xD5)
C_LLM         = RGBColor(0x05, 0x96, 0x69)   # LLM green

# ---------------------------------------------------------------------------
# Helper utilities
# ---------------------------------------------------------------------------

def rgb(r, g, b):
    return RGBColor(r, g, b)


def add_rect(slide, x, y, w, h, fill=None, line_color=None, line_width=Pt(0.75), radius=None):
    """Add a rectangle shape and return it."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill

    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = line_width

    return shape


def add_text_box(slide, text, x, y, w, h, font_size=Pt(10), bold=False,
                  color=None, align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return txBox


def add_labeled_rect(slide, x, y, w, h, label, fill, line_color=None,
                      font_size=Pt(9), bold=False, text_color=None):
    shape = add_rect(slide, x, y, w, h, fill=fill, line_color=line_color)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = text_color or C_TITLE_TEXT
    return shape


def add_connector(slide, x1, y1, x2, y2):
    """Add a simple straight line connector."""
    from pptx.util import Inches as I
    connector = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR_TYPE.STRAIGHT
        I(x1), I(y1), I(x2), I(y2)
    )
    connector.line.color.rgb = C_ARROW
    connector.line.width = Pt(1)
    return connector


# ---------------------------------------------------------------------------
# Build the presentation
# ---------------------------------------------------------------------------

def build_pptx(output_path: str):
    prs = Presentation()

    # Widescreen 16:9 slide
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Slide background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = C_BG

    # -----------------------------------------------------------------------
    # Title bar
    # -----------------------------------------------------------------------
    add_rect(slide, 0, 0, 13.33, 0.55, fill=C_HEADER)
    add_text_box(slide, "SmartRecover — Agentic Incident Management Platform",
                 0.15, 0.08, 10, 0.4,
                 font_size=Pt(16), bold=True, color=C_HEADER_TEXT, align=PP_ALIGN.LEFT)
    add_text_box(slide, "Architecture Overview",
                 10.15, 0.13, 3, 0.3,
                 font_size=Pt(10), color=RGBColor(0xC4, 0xB5, 0xFD), align=PP_ALIGN.RIGHT)

    # -----------------------------------------------------------------------
    # Section labels (row below header)
    # -----------------------------------------------------------------------
    section_y = 0.6
    labels = [
        (0.1,  3.0, "① Incident List (ServiceNow-style)"),
        (3.3,  5.1, "② Ticket Details + Agent Results"),
        (8.6,  4.55, "③ Streaming Chat"),
    ]
    for lx, lw, ltxt in labels:
        add_text_box(slide, ltxt, lx, section_y, lw, 0.25,
                     font_size=Pt(8), bold=True, color=C_SUBTEXT, align=PP_ALIGN.LEFT)

    # -----------------------------------------------------------------------
    # Panel separators (vertical lines)
    # -----------------------------------------------------------------------
    for sep_x in (3.2, 8.5):
        add_connector(slide, sep_x, 0.55, sep_x, 7.45)

    # -----------------------------------------------------------------------
    # LEFT PANEL — Incident List
    # -----------------------------------------------------------------------
    panel_y = 0.9

    # Panel background
    add_rect(slide, 0.05, panel_y - 0.05, 3.1, 6.55,
             fill=C_PANEL_BG, line_color=C_BORDER)

    # Filter buttons row
    filter_buttons = [("All", C_AGENT, C_HEADER_TEXT), ("Open", C_AGENT_LIGHT, C_AGENT),
                      ("Investigating", C_AGENT_LIGHT, C_AGENT), ("Resolved", C_AGENT_LIGHT, C_AGENT)]
    fx = 0.12
    for fl, fbg, ftxt in filter_buttons:
        bw = 0.56 if fl == "Investigating" else 0.46
        add_labeled_rect(slide, fx, panel_y, bw, 0.22, fl, fill=fbg,
                          font_size=Pt(7), text_color=ftxt)
        fx += bw + 0.05

    # Incident cards (ServiceNow style)
    incidents = [
        ("INC001", "Memory leak in auth service",       C_MED,  C_YELLOW_BG, "3 - Moderate", "In Progress", "search-team",    "Jan 17, 2026"),
        ("INC002", "API response latency spike",        C_MED,  C_YELLOW_BG, "3 - Moderate", "New",         "dba-team",       "Jan 15, 2026"),
        ("INC003", "Log aggregation pipeline broken",   C_HIGH, C_ORANGE_BG, "2 - High",     "New",         "payments-team",  "Jan 19, 2026"),
        ("INC004", "Redis cache connection failures",   C_LOW,  C_BLUE_BG,   "4 - Low",      "New",         "ops-team",       "Jan 12, 2026"),
        ("INC005", "Rate limiting misconfiguration",    C_LOW,  C_BLUE_BG,   "4 - Low",      "Resolved",    "api-team",       "Jan 11, 2026"),
    ]

    card_y = panel_y + 0.28
    for inc_id, inc_title, bar_color, pri_bg, pri_label, state, assignee, date in incidents:
        card_h = 0.78
        # card body
        add_rect(slide, 0.1, card_y, 3.0, card_h,
                 fill=C_PANEL_BG, line_color=C_BORDER)
        # left priority bar
        add_rect(slide, 0.1, card_y, 0.06, card_h, fill=bar_color)
        # incident number
        add_text_box(slide, inc_id, 0.22, card_y + 0.03, 0.7, 0.16,
                     font_size=Pt(7.5), bold=True, color=C_AGENT)
        # state badge
        state_color = C_GREEN if state == "Resolved" else (C_MED if state == "In Progress" else C_LOW)
        state_bg    = C_GREEN_BG if state == "Resolved" else (C_YELLOW_BG if state == "In Progress" else C_BLUE_BG)
        add_labeled_rect(slide, 2.35, card_y + 0.04, 0.68, 0.15, state,
                          fill=state_bg, font_size=Pt(6.5), text_color=state_color)
        # short description
        add_text_box(slide, inc_title, 0.22, card_y + 0.19, 2.8, 0.22,
                     font_size=Pt(8), bold=False, color=C_TITLE_TEXT, wrap=True)
        # priority badge
        add_labeled_rect(slide, 0.22, card_y + 0.41, 0.75, 0.14, pri_label,
                          fill=pri_bg, font_size=Pt(6), text_color=bar_color)
        # assignee + date
        add_text_box(slide, f"👤 {assignee}", 0.22, card_y + 0.57, 1.5, 0.14,
                     font_size=Pt(6.5), color=C_SUBTEXT)
        add_text_box(slide, f"📅 {date}", 1.78, card_y + 0.57, 1.25, 0.14,
                     font_size=Pt(6.5), color=C_SUBTEXT, align=PP_ALIGN.RIGHT)
        card_y += card_h + 0.05

    # "... X more incidents" hint
    add_text_box(slide, "⋯  25 more incidents", 0.1, card_y + 0.05, 3.0, 0.2,
                 font_size=Pt(7.5), italic=True, color=C_SUBTEXT, align=PP_ALIGN.CENTER)

    # -----------------------------------------------------------------------
    # MIDDLE PANEL — Ticket Details + Agent Results
    # -----------------------------------------------------------------------
    mid_x = 3.3
    add_rect(slide, mid_x, panel_y - 0.05, 5.05, 6.55,
             fill=C_PANEL_BG, line_color=C_BORDER)

    # Incident header
    add_rect(slide, mid_x, panel_y - 0.05, 5.05, 0.5,
             fill=RGBColor(0xF5, 0xF3, 0xFF), line_color=C_BORDER)
    add_text_box(slide, "INC001 — Memory leak in auth service",
                 mid_x + 0.1, panel_y, 4.0, 0.25,
                 font_size=Pt(10), bold=True, color=C_TITLE_TEXT)
    add_text_box(slide, "Severity: Medium  |  Assignee: search-team  |  Opened: Jan 17, 2026",
                 mid_x + 0.1, panel_y + 0.25, 4.5, 0.18,
                 font_size=Pt(7.5), color=C_SUBTEXT)

    # Agent result tabs
    tabs = ["ServiceNow", "Knowledge Base", "Changes", "Logs", "Events", "Remediations"]
    tab_x = mid_x + 0.1
    tab_y = panel_y + 0.5
    tab_colors = [
        (C_AGENT, C_HEADER_TEXT),  # active
        (C_BORDER, C_SUBTEXT),
        (C_BORDER, C_SUBTEXT),
        (C_BORDER, C_SUBTEXT),
        (C_BORDER, C_SUBTEXT),
        (C_BORDER, C_SUBTEXT),
    ]
    for i, (tab, (tbg, tfg)) in enumerate(zip(tabs, tab_colors)):
        tw = 4.8 / len(tabs)
        add_labeled_rect(slide, tab_x + i * tw, tab_y, tw - 0.02, 0.22, tab,
                          fill=tbg, font_size=Pt(7), text_color=tfg, bold=(i == 0))

    # ServiceNow results panel
    res_y = tab_y + 0.27
    add_rect(slide, mid_x + 0.08, res_y, 4.86, 3.5,
             fill=RGBColor(0xFA, 0xFA, 0xFF), line_color=C_BORDER)
    add_text_box(slide, "Similar Incidents",
                 mid_x + 0.18, res_y + 0.08, 3, 0.18,
                 font_size=Pt(8), bold=True, color=C_TITLE_TEXT)

    sim_incs = [
        ("INC042", "Memory leak in reporting service", "92% match", C_GREEN, C_GREEN_BG),
        ("INC017", "Auth service OOM crash",           "84% match", C_MED,   C_YELLOW_BG),
        ("INC009", "Cache eviction anomaly",           "71% match", C_LOW,   C_BLUE_BG),
    ]
    siy = res_y + 0.3
    for sid, stitle, score, sc, sbg in sim_incs:
        add_rect(slide, mid_x + 0.12, siy, 4.75, 0.55,
                 fill=C_PANEL_BG, line_color=C_BORDER)
        add_text_box(slide, sid, mid_x + 0.2, siy + 0.06, 0.55, 0.16,
                     font_size=Pt(7.5), bold=True, color=C_AGENT)
        add_text_box(slide, stitle, mid_x + 0.78, siy + 0.06, 3.0, 0.16,
                     font_size=Pt(7.5), color=C_TITLE_TEXT)
        add_labeled_rect(slide, mid_x + 3.85, siy + 0.06, 0.8, 0.16, score,
                          fill=sbg, font_size=Pt(6.5), text_color=sc)
        add_text_box(slide, "Resolution: Restart auth pods; increase heap to 2 GB",
                     mid_x + 0.2, siy + 0.28, 4.3, 0.16,
                     font_size=Pt(6.5), italic=True, color=C_SUBTEXT)
        siy += 0.62

    # Orchestrator / retrieve context button
    add_rect(slide, mid_x + 0.12, res_y + 3.6, 2.1, 0.3,
             fill=C_AGENT, line_color=None)
    add_text_box(slide, "⚡  Retrieve Context", mid_x + 0.22, res_y + 3.66, 1.9, 0.2,
                 font_size=Pt(8.5), bold=True, color=C_HEADER_TEXT, align=PP_ALIGN.CENTER)

    add_rect(slide, mid_x + 2.4, res_y + 3.6, 2.45, 0.3,
             fill=RGBColor(0xF3, 0xF4, 0xF6), line_color=C_BORDER)
    add_text_box(slide, "✗  Exclude Selected Items", mid_x + 2.5, res_y + 3.66, 2.25, 0.2,
                 font_size=Pt(8), color=C_SUBTEXT, align=PP_ALIGN.CENTER)

    # -----------------------------------------------------------------------
    # RIGHT PANEL — Chat
    # -----------------------------------------------------------------------
    chat_x = 8.6
    add_rect(slide, chat_x, panel_y - 0.05, 4.6, 6.55,
             fill=C_PANEL_BG, line_color=C_BORDER)

    add_text_box(slide, "AI Assistant",
                 chat_x + 0.15, panel_y, 3, 0.25,
                 font_size=Pt(10), bold=True, color=C_TITLE_TEXT)
    add_text_box(slide, "Streaming responses via LangGraph",
                 chat_x + 0.15, panel_y + 0.24, 4.2, 0.16,
                 font_size=Pt(7), italic=True, color=C_SUBTEXT)

    # Chat bubbles
    bubbles = [
        ("What caused this memory leak?",
         "user",
         RGBColor(0xED, 0xE9, 0xFE), C_TITLE_TEXT),
        ("Based on similar incident INC042 and the recent database schema change "
         "(CHG016 on Jan 7), the root cause is likely an unreleased JDBC connection "
         "pool. Recommend restarting auth pods and patching the connection pool config.",
         "ai",
         RGBColor(0xF0, 0xFD, 0xFA), C_TITLE_TEXT),
        ("Which runbook should I follow?",
         "user",
         RGBColor(0xED, 0xE9, 0xFE), C_TITLE_TEXT),
    ]

    bub_y = panel_y + 0.5
    for text, role, bbg, btxt in bubbles:
        is_user = role == "user"
        bx = chat_x + (1.2 if is_user else 0.1)
        bw = 3.3
        bh = 0.55 if is_user else 1.0
        add_rect(slide, bx, bub_y, bw, bh, fill=bbg, line_color=C_BORDER)
        add_text_box(slide, text, bx + 0.08, bub_y + 0.06, bw - 0.15, bh - 0.1,
                     font_size=Pt(7.5), color=btxt, wrap=True,
                     align=PP_ALIGN.RIGHT if is_user else PP_ALIGN.LEFT)
        bub_y += bh + 0.1

    # Streaming indicator
    add_text_box(slide, "● Streaming…",
                 chat_x + 0.15, bub_y + 0.05, 2, 0.2,
                 font_size=Pt(7.5), italic=True, color=C_AGENT)

    # Chat input box
    add_rect(slide, chat_x + 0.1, 6.9, 3.6, 0.3,
             fill=C_PANEL_BG, line_color=C_BORDER)
    add_text_box(slide, "Ask a follow-up question…",
                 chat_x + 0.2, 6.95, 3.2, 0.2,
                 font_size=Pt(8), italic=True, color=C_SUBTEXT)
    add_rect(slide, chat_x + 3.8, 6.9, 0.65, 0.3, fill=C_AGENT)
    add_text_box(slide, "Send", chat_x + 3.82, 6.95, 0.6, 0.2,
                 font_size=Pt(8), bold=True, color=C_HEADER_TEXT, align=PP_ALIGN.CENTER)

    # -----------------------------------------------------------------------
    # BOTTOM — LangGraph Orchestration legend
    # -----------------------------------------------------------------------
    orch_y = 7.08
    add_rect(slide, 0.05, orch_y - 0.05, 13.23, 0.42,
             fill=RGBColor(0xF5, 0xF3, 0xFF), line_color=C_BORDER)

    add_text_box(slide, "LangGraph Orchestration:",
                 0.15, orch_y, 1.5, 0.25,
                 font_size=Pt(7.5), bold=True, color=C_AGENT)

    agents = [
        ("Orchestrator", C_AGENT,              C_HEADER_TEXT),
        ("ServiceNow Agent", C_AGENT_LIGHT,    C_AGENT),
        ("Knowledge Base Agent", C_AGENT_LIGHT, C_AGENT),
        ("Change Correlation", C_AGENT_LIGHT,  C_AGENT),
        ("Logs Agent", C_AGENT_LIGHT,          C_AGENT),
        ("Events Agent", C_AGENT_LIGHT,        C_AGENT),
        ("LLM (OpenAI / Gemini / Ollama)", RGBColor(0xD1, 0xFA, 0xE5), C_LLM),
    ]
    ax = 1.75
    for aname, abg, afg in agents:
        aw = len(aname) * 0.075 + 0.25
        add_labeled_rect(slide, ax, orch_y + 0.03, aw, 0.22, aname,
                          fill=abg, font_size=Pt(7), text_color=afg, bold=(aname == "Orchestrator"))
        ax += aw + 0.08
        if ax < 12.5:
            add_text_box(slide, "→", ax - 0.1, orch_y + 0.05, 0.12, 0.18,
                         font_size=Pt(8), color=C_ARROW, align=PP_ALIGN.CENTER)

    # -----------------------------------------------------------------------
    # Save
    # -----------------------------------------------------------------------
    prs.save(output_path)
    print(f"Saved: {output_path}")


if __name__ == "__main__":
    script_dir = os.path.dirname(os.path.abspath(__file__))
    output = os.path.join(script_dir, "SmartRecover_Architecture.pptx")
    build_pptx(output)
